#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Universal Document Parser v2.0 (N合一文档解析器)
改进版：懒加载、OCR 内存优化、大文件截断

支持格式：Word(.doc/.docx/.wps)、Excel(.xls/.xlsx/.xlsm/.et)、
         PowerPoint(.ppt/.pptx/.dps)、PDF(电子+扫描件OCR)、
         纯文本、图片OCR
"""

import os
import sys
import json
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Union, Optional, List, Dict, Any, Generator
from dataclasses import dataclass, field


# ============================================================
# 懒加载机制：用到才 import，缺库就报错但不阻塞启动
# ============================================================

_LAZY_MODULES = {}

def _lazy_import(name: str):
    """懒加载模块，第一次使用时才 import"""
    if name not in _LAZY_MODULES:
        try:
            if name == "docx":
                import docx as _mod
                _LAZY_MODULES[name] = _mod
            elif name == "pptx":
                import pptx as _mod
                _LAZY_MODULES[name] = _mod
            elif name == "openpyxl":
                import openpyxl as _mod
                _LAZY_MODULES[name] = _mod
            elif name == "xlrd":
                import xlrd as _mod
                _LAZY_MODULES[name] = _mod
            elif name == "pdfplumber":
                import pdfplumber as _mod
                _LAZY_MODULES[name] = _mod
            elif name == "fitz":
                import fitz as _mod
                _LAZY_MODULES[name] = _mod
            elif name == "pytesseract":
                import pytesseract as _mod
                _LAZY_MODULES[name] = _mod
            elif name == "PIL":
                from PIL import Image as _mod
                _LAZY_MODULES[name] = _mod
            elif name == "mammoth":
                import mammoth as _mod
                _LAZY_MODULES[name] = _mod
            elif name == "pdf2image":
                from pdf2image import convert_from_path as _mod
                _LAZY_MODULES[name] = _mod
            elif name == "pandas":
                import pandas as _mod
                _LAZY_MODULES[name] = _mod
        except ImportError:
            _LAZY_MODULES[name] = None
    return _LAZY_MODULES[name]


def _check_dep(name: str, feature: str) -> None:
    """检查依赖是否存在，不存在则抛出清晰错误"""
    mod = _lazy_import(name)
    if mod is None:
        pip_name = "Pillow" if name == "PIL" else name
        raise ImportError(f"缺少依赖库 '{name}'，无法使用 {feature} 功能。请运行: pip install {pip_name}")


# ============================================================
# 配置常量
# ============================================================

MAX_OCR_PAGES = 50          # PDF OCR 最大页数，防止内存爆炸
MAX_TEXT_LENGTH = 15000     # 返回文本最大长度，超出则截断
OCR_DPI = 200               # OCR 分辨率，降低以节省内存
MAX_FILE_SIZE = 50 * 1024 * 1024   # 50MB：超大文件直接拒绝，防解析器内存爆炸


# ============================================================
# 数据结构
# ============================================================

@dataclass
class ParseResult:
    """解析结果容器"""
    text: str = ""
    metadata: Dict[str, Any] = field(default_factory=dict)
    pages: List[str] = field(default_factory=list)
    tables: List[List[List[str]]] = field(default_factory=list)
    success: bool = True
    error: str = ""
    method: str = ""
    truncated: bool = False   # 是否被截断

    def to_dict(self) -> dict:
        return {
            "text": self.text,
            "metadata": self.metadata,
            "pages": self.pages,
            "tables_count": len(self.tables),
            "success": self.success,
            "error": self.error,
            "method": self.method,
            "truncated": self.truncated,
        }

    def __repr__(self) -> str:
        preview = self.text[:200].replace("\n", " ") if self.text else "(空)"
        return f"ParseResult(success={self.success}, method={self.method}, truncated={self.truncated}, text_preview={preview!r})"


# ============================================================
# 工具函数
# ============================================================

def _find_libreoffice() -> Optional[str]:
    """查找 LibreOffice / soffice 可执行文件"""
    for cmd in ["soffice", "libreoffice"]:
        path = shutil.which(cmd)
        if path:
            return path
    fallback_paths = [
        "/usr/bin/soffice",
        "/usr/bin/libreoffice",
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "C:/Program Files/LibreOffice/program/soffice.exe",
        "C:/Program Files (x86)/LibreOffice/program/soffice.exe",
    ]
    for p in fallback_paths:
        if os.path.exists(p):
            return p
    return None


def _libreoffice_convert(src: Path, target_ext: str, out_dir: Path) -> Optional[Path]:
    """用 LibreOffice 将文件转换为 target_ext 格式"""
    soffice = _find_libreoffice()
    if not soffice:
        return None
    cmd = [
        soffice, "--headless", "--convert-to", target_ext,
        "--outdir", str(out_dir), str(src)
    ]
    try:
        subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=60)
        converted = out_dir / (src.stem + "." + target_ext)
        if converted.exists():
            return converted
    except Exception:
        pass
    return None


def _is_scanned_pdf(file_path: Path) -> bool:
    """启发式判断 PDF 是否为扫描件（图片型）"""
    fitz = _lazy_import("fitz")
    if not fitz:
        return False
    try:
        doc = fitz.open(str(file_path))
        total_pages = len(doc)
        text_len = 0
        image_count = 0
        for page in doc:
            text_len += len(page.get_text().strip())
            image_count += len(page.get_images())
        doc.close()
        avg_text = text_len / max(total_pages, 1)
        if avg_text < 50 and image_count > 0:
            return True
        return False
    except Exception:
        return False


def _ocr_pdf_pages(file_path: Path, lang: str = "chi_sim+eng") -> Generator[str, None, None]:
    """逐页 OCR PDF，生成器模式，避免内存爆炸"""
    _check_dep("pdf2image", "PDF 转图片")
    _check_dep("pytesseract", "OCR 文字识别")
    _check_dep("PIL", "图片处理")

    convert_from_path = _lazy_import("pdf2image")
    pytesseract = _lazy_import("pytesseract")

    total_pages = 0
    try:
        fitz = _lazy_import("fitz")
        if fitz:
            doc = fitz.open(str(file_path))
            total_pages = len(doc)
            doc.close()
    except Exception:
        pass

    for page_num in range(1, min(total_pages if total_pages else MAX_OCR_PAGES, MAX_OCR_PAGES) + 1):
        try:
            images = convert_from_path(str(file_path), dpi=OCR_DPI, first_page=page_num, last_page=page_num)
            if images:
                text = pytesseract.image_to_string(images[0], lang=lang)
                yield f"--- 第 {page_num} 页 ---\n{text}"
            del images
        except Exception as e:
            yield f"--- 第 {page_num} 页 [OCR 失败: {e}] ---"

    if total_pages > MAX_OCR_PAGES:
        yield f"\n...（PDF 共 {total_pages} 页，已截断至前 {MAX_OCR_PAGES} 页）"


def _ocr_pdf(file_path: Path, lang: str = "chi_sim+eng") -> ParseResult:
    """对 PDF 逐页 OCR（生成器模式，内存安全）"""
    result = ParseResult(method="pdf_ocr")
    try:
        pages_text = list(_ocr_pdf_pages(file_path, lang))
        result.pages = pages_text
        result.text = "\n\n".join(pages_text)
        result.metadata["ocr_engine"] = "tesseract"
        result.metadata["ocr_dpi"] = OCR_DPI
        result.metadata["max_pages"] = MAX_OCR_PAGES
    except ImportError as e:
        result.success = False
        result.error = str(e)
    except Exception as e:
        result.success = False
        result.error = f"PDF OCR 失败: {e}"
    return result


def _ocr_image(file_path: Path, lang: str = "chi_sim+eng") -> ParseResult:
    """对图片 OCR"""
    result = ParseResult(method="image_ocr")
    _check_dep("pytesseract", "图片 OCR")
    _check_dep("PIL", "图片处理")
    try:
        Image = _lazy_import("PIL")
        pytesseract = _lazy_import("pytesseract")
        img = Image.open(str(file_path))
        text = pytesseract.image_to_string(img, lang=lang)
        result.text = text
        result.metadata["image_size"] = img.size
        result.metadata["ocr_engine"] = "tesseract"
    except Exception as e:
        result.success = False
        result.error = f"图片 OCR 失败: {e}"
    return result


def _truncate_text(result: ParseResult) -> ParseResult:
    """大文件截断：如果文本超过 MAX_TEXT_LENGTH，截断并标记"""
    original_length = len(result.text)
    if original_length > MAX_TEXT_LENGTH:
        result.text = result.text[:MAX_TEXT_LENGTH] + f"\n\n...（文档共 {original_length} 字，已截断至前 {MAX_TEXT_LENGTH} 字）"
        result.truncated = True
        result.metadata["original_length"] = original_length
        result.metadata["truncated_to"] = MAX_TEXT_LENGTH
    return result


# ============================================================
# 各格式解析器
# ============================================================

def parse_docx(file_path: Path) -> ParseResult:
    """解析 Word .docx"""
    result = ParseResult(method="python-docx")
    _check_dep("docx", "Word .docx 解析")
    try:
        docx = _lazy_import("docx")
        doc = docx.Document(str(file_path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        tables = []
        for table in doc.tables:
            tbl = []
            for row in table.rows:
                tbl.append([cell.text for cell in row.cells])
            tables.append(tbl)
        result.text = "\n".join(paragraphs)
        result.pages = paragraphs
        result.tables = tables
        result.metadata["paragraphs"] = len(paragraphs)
        result.metadata["tables"] = len(tables)
    except Exception as e:
        result.success = False
        result.error = f"docx 解析失败: {e}"
    return _truncate_text(result)


def parse_doc(file_path: Path) -> ParseResult:
    """解析 Word .doc（旧版）"""
    mammoth = _lazy_import("mammoth")
    if mammoth:
        try:
            with open(file_path, "rb") as f:
                res = mammoth.extract_raw_text(f)
            return _truncate_text(ParseResult(
                text=res.value,
                method="mammoth",
                metadata={"mammoth_messages": len(res.messages)}
            ))
        except Exception:
            pass
    with tempfile.TemporaryDirectory() as tmpdir:
        converted = _libreoffice_convert(file_path, "docx", Path(tmpdir))
        if converted:
            res = parse_docx(converted)
            res.method = "libreoffice->docx"
            return _truncate_text(res)
    antiword = shutil.which("antiword")
    if antiword:
        try:
            out = subprocess.check_output([antiword, str(file_path)], stderr=subprocess.DEVNULL, text=True)
            return _truncate_text(ParseResult(text=out, method="antiword"))
        except Exception:
            pass
    return ParseResult(success=False, error="无法解析 .doc 文件，请安装 mammoth、LibreOffice 或 antiword")


def parse_wps(file_path: Path) -> ParseResult:
    """解析 WPS 文字 .wps"""
    with tempfile.TemporaryDirectory() as tmpdir:
        converted = _libreoffice_convert(file_path, "docx", Path(tmpdir))
        if converted:
            res = parse_docx(converted)
            res.method = "libreoffice(wps)->docx"
            return _truncate_text(res)
    return parse_doc(file_path)


def parse_xlsx(file_path: Path) -> ParseResult:
    """解析 Excel .xlsx / .xlsm"""
    result = ParseResult(method="openpyxl")
    _check_dep("openpyxl", "Excel .xlsx 解析")
    try:
        openpyxl = _lazy_import("openpyxl")
        wb = openpyxl.load_workbook(str(file_path), data_only=True)
        sheets_text = []
        all_tables = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                row_text = [str(cell) if cell is not None else "" for cell in row]
                rows.append(row_text)
            rows = [r for r in rows if any(c.strip() for c in r)]
            if rows:
                sheet_text = f"--- Sheet: {sheet_name} ---\n"
                sheet_text += "\n".join(["\t".join(r) for r in rows])
                sheets_text.append(sheet_text)
                all_tables.append(rows)
        result.text = "\n\n".join(sheets_text)
        result.pages = sheets_text
        result.tables = all_tables
        result.metadata["sheets"] = len(wb.sheetnames)
    except Exception as e:
        result.success = False
        result.error = f"xlsx 解析失败: {e}"
    return _truncate_text(result)


def parse_xls(file_path: Path) -> ParseResult:
    """解析 Excel .xls（旧版）"""
    xlrd = _lazy_import("xlrd")
    if not xlrd:
        with tempfile.TemporaryDirectory() as tmpdir:
            converted = _libreoffice_convert(file_path, "xlsx", Path(tmpdir))
            if converted:
                res = parse_xlsx(converted)
                res.method = "libreoffice(xls)->xlsx"
                return _truncate_text(res)
        return ParseResult(success=False, error="缺少 xlrd 库，且 LibreOffice 转换失败")
    result = ParseResult(method="xlrd")
    try:
        wb = xlrd.open_workbook(str(file_path))
        sheets_text = []
        all_tables = []
        for sheet_idx in range(wb.nsheets):
            ws = wb.sheet_by_index(sheet_idx)
            rows = []
            for r in range(ws.nrows):
                row = [str(ws.cell_value(r, c)) for c in range(ws.ncols)]
                rows.append(row)
            rows = [r for r in rows if any(c.strip() for c in r)]
            if rows:
                sheet_text = f"--- Sheet: {ws.name} ---\n"
                sheet_text += "\n".join(["\t".join(r) for r in rows])
                sheets_text.append(sheet_text)
                all_tables.append(rows)
        result.text = "\n\n".join(sheets_text)
        result.pages = sheets_text
        result.tables = all_tables
        result.metadata["sheets"] = wb.nsheets
    except Exception as e:
        result.success = False
        result.error = f"xls 解析失败: {e}"
    return _truncate_text(result)


def parse_et(file_path: Path) -> ParseResult:
    """解析 WPS 表格 .et"""
    with tempfile.TemporaryDirectory() as tmpdir:
        converted = _libreoffice_convert(file_path, "xlsx", Path(tmpdir))
        if converted:
            res = parse_xlsx(converted)
            res.method = "libreoffice(et)->xlsx"
            return _truncate_text(res)
    return parse_xls(file_path)


def parse_pptx(file_path: Path) -> ParseResult:
    """解析 PowerPoint .pptx"""
    result = ParseResult(method="python-pptx")
    _check_dep("pptx", "PowerPoint .pptx 解析")
    try:
        pptx = _lazy_import("pptx")
        prs = pptx.Presentation(str(file_path))
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            slide_lines = [f"--- 第 {i} 页 ---"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_lines.append(shape.text.strip())
                if shape.has_table:
                    tbl = []
                    for row in shape.table.rows:
                        tbl.append([cell.text for cell in row.cells])
                    if tbl:
                        slide_lines.append("[表格]\n" + "\n".join(["\t".join(r) for r in tbl]))
                        result.tables.append(tbl)
            if len(slide_lines) > 1:
                slides_text.append("\n".join(slide_lines))
        result.text = "\n\n".join(slides_text)
        result.pages = slides_text
        result.metadata["slides"] = len(prs.slides)
    except Exception as e:
        result.success = False
        result.error = f"pptx 解析失败: {e}"
    return _truncate_text(result)


def parse_ppt(file_path: Path) -> ParseResult:
    """解析 PowerPoint .ppt（旧版）"""
    with tempfile.TemporaryDirectory() as tmpdir:
        converted = _libreoffice_convert(file_path, "pptx", Path(tmpdir))
        if converted:
            res = parse_pptx(converted)
            res.method = "libreoffice(ppt)->pptx"
            return _truncate_text(res)
    return ParseResult(success=False, error="无法解析 .ppt 文件，请安装 LibreOffice")


def parse_dps(file_path: Path) -> ParseResult:
    """解析 WPS 演示 .dps"""
    with tempfile.TemporaryDirectory() as tmpdir:
        converted = _libreoffice_convert(file_path, "pptx", Path(tmpdir))
        if converted:
            res = parse_pptx(converted)
            res.method = "libreoffice(dps)->pptx"
            return _truncate_text(res)
    return parse_ppt(file_path)


def parse_pdf_text(file_path: Path) -> ParseResult:
    """解析 PDF（文本提取模式）"""
    result = ParseResult(method="pdfplumber")
    pdfplumber = _lazy_import("pdfplumber")
    if pdfplumber:
        try:
            pages_text = []
            all_tables = []
            with pdfplumber.open(str(file_path)) as pdf:
                for i, page in enumerate(pdf.pages, 1):
                    text = page.extract_text() or ""
                    if text.strip():
                        pages_text.append(f"--- 第 {i} 页 ---\n{text.strip()}")
                    tables = page.extract_tables()
                    for tbl in tables:
                        if tbl:
                            all_tables.append(tbl)
            result.text = "\n\n".join(pages_text)
            result.pages = pages_text
            result.tables = all_tables
            result.metadata["pages"] = len(pages_text)
            result.metadata["tables"] = len(all_tables)
            return _truncate_text(result)
        except Exception as e:
            result.error = f"pdfplumber 失败: {e}"
    fitz = _lazy_import("fitz")
    if fitz:
        try:
            doc = fitz.open(str(file_path))
            pages_text = []
            for i, page in enumerate(doc, 1):
                text = page.get_text().strip()
                if text:
                    pages_text.append(f"--- 第 {i} 页 ---\n{text}")
            doc.close()
            result.text = "\n\n".join(pages_text)
            result.pages = pages_text
            result.metadata["pages"] = len(pages_text)
            result.method = "pymupdf"
            return _truncate_text(result)
        except Exception as e:
            result.error += f" | pymupdf 失败: {e}"
    result.success = False
    result.error = result.error or "无法解析 PDF，请安装 pdfplumber 或 pymupdf"
    return result


def parse_pdf(file_path: Path, force_ocr: bool = False) -> ParseResult:
    """解析 PDF（自动判断是否需要 OCR）"""
    if force_ocr:
        return _ocr_pdf(file_path)
    text_res = parse_pdf_text(file_path)
    if text_res.success and text_res.text.strip():
        if len(text_res.text.strip()) < 200 and _is_scanned_pdf(file_path):
            return _ocr_pdf(file_path)
        return text_res
    return _ocr_pdf(file_path)


# ============================================================
# 主入口：万能解析器
# ============================================================

PARSER_MAP = {
    ".docx": parse_docx,
    ".doc": parse_doc,
    ".wps": parse_wps,
    ".xlsx": parse_xlsx,
    ".xlsm": parse_xlsx,
    ".xls": parse_xls,
    ".et": parse_et,
    ".pptx": parse_pptx,
    ".ppt": parse_ppt,
    ".dps": parse_dps,
    ".pdf": parse_pdf,
    ".png": _ocr_image,
    ".jpg": _ocr_image,
    ".jpeg": _ocr_image,
    ".bmp": _ocr_image,
    ".tiff": _ocr_image,
    ".gif": _ocr_image,
    ".webp": _ocr_image,
}


def parse_document(file_path: Union[str, Path], force_ocr: bool = False, lang: str = "chi_sim+eng") -> ParseResult:
    """
    万能文档解析入口

    Args:
        file_path: 文件路径
        force_ocr: 对 PDF 强制使用 OCR
        lang: OCR 语言，默认中文简体+英文

    Returns:
        ParseResult 对象
    """
    path = Path(file_path)
    if not path.exists():
        return ParseResult(success=False, error=f"文件不存在: {file_path}")

    # 大文件保护：超过上限直接拒绝（python-docx/pdfplumber 等会整体载入内存）
    try:
        size = path.stat().st_size
    except OSError:
        size = 0
    if size > MAX_FILE_SIZE:
        return ParseResult(
            success=False,
            error=(f"文件过大（{size / 1024 / 1024:.1f}MB > "
                   f"{MAX_FILE_SIZE // (1024 * 1024)}MB），请手动处理或拆分后重试"))

    ext = path.suffix.lower()

    if ext in (".txt", ".md", ".csv", ".json", ".xml", ".html", ".htm", ".py", ".js", ".css", ".yaml", ".yml"):
        try:
            text = path.read_text(encoding="utf-8", errors="ignore")
            return _truncate_text(ParseResult(
                text=text,
                method="direct_text",
                metadata={"encoding": "utf-8", "size": path.stat().st_size}
            ))
        except Exception as e:
            return ParseResult(success=False, error=f"文本读取失败: {e}")

    parser = PARSER_MAP.get(ext)
    if parser:
        if ext == ".pdf":
            return parser(path, force_ocr=force_ocr)
        return parser(path)

    try:
        text = path.read_text(encoding="utf-8", errors="ignore")
        if text.strip():
            return _truncate_text(ParseResult(text=text, method="fallback_text"))
    except Exception:
        pass

    return ParseResult(
        success=False,
        error=f"不支持的文件格式: {ext}。支持的格式: {', '.join(PARSER_MAP.keys())} 及常见纯文本格式"
    )


def batch_parse(folder: Union[str, Path], pattern: str = "*", recursive: bool = False) -> Dict[str, ParseResult]:
    """批量解析文件夹内的文档"""
    folder = Path(folder)
    if not folder.is_dir():
        return {}
    results = {}
    iterator = folder.rglob(pattern) if recursive else folder.glob(pattern)
    for fp in iterator:
        if fp.is_file() and fp.suffix.lower() in PARSER_MAP:
            results[str(fp)] = parse_document(fp)
    return results


# ============================================================
# CLI 命令行接口
# ============================================================

if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="Universal Document Parser v2.0 (N合一文档解析器)")
    parser.add_argument("file", help="要解析的文件路径")
    parser.add_argument("-o", "--output", help="输出文件路径（默认 stdout）")
    parser.add_argument("--force-ocr", action="store_true", help="PDF 强制 OCR")
    parser.add_argument("--lang", default="chi_sim+eng", help="OCR 语言（默认 chi_sim+eng）")
    parser.add_argument("--json", action="store_true", help="以 JSON 格式输出")
    parser.add_argument("--preview", type=int, default=500, help="文本预览长度（默认 500）")

    args = parser.parse_args()

    try:
        res = parse_document(args.file, force_ocr=args.force_ocr, lang=args.lang)
    except ImportError as e:
        print(f"❌ 依赖错误: {e}", file=sys.stderr)
        sys.exit(1)
    except Exception as e:
        print(f"❌ 解析失败: {e}", file=sys.stderr)
        sys.exit(1)

    if args.json:
        output = json.dumps(res.to_dict(), ensure_ascii=False, indent=2)
    else:
        lines = [
            f"📄 文件: {args.file}",
            f"✅ 成功: {res.success}",
            f"🔧 方法: {res.method}",
            f"✂️ 截断: {res.truncated}",
            f"📊 元数据: {json.dumps(res.metadata, ensure_ascii=False)}",
            f"📋 文本预览（前 {args.preview} 字）:",
            "-" * 50,
            res.text[:args.preview] + ("..." if len(res.text) > args.preview else ""),
            "-" * 50,
        ]
        if res.error:
            lines.append(f"❌ 错误: {res.error}")
        output = "\n".join(lines)

    if args.output:
        Path(args.output).write_text(output, encoding="utf-8")
        print(f"结果已保存到: {args.output}")
    else:
        print(output)
